#!/usr/bin/env python3
"""
《天朝的崩溃：鸦片战争再研究》精美PPT生成器
深蓝灰 + 暖米色 + 古铜金 配色方案
每页背景根据内容定制，知识密度高
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ===== 配色 =====
NAVY       = RGBColor(0x1B, 0x2A, 0x4A)
DARK_NAVY  = RGBColor(0x0F, 0x1A, 0x33)
DEEP_NAVY  = RGBColor(0x0A, 0x12, 0x25)
CREAM      = RGBColor(0xF5, 0xF0, 0xE8)
WARM_CREAM = RGBColor(0xED, 0xE5, 0xD8)
PARCHMENT  = RGBColor(0xE8, 0xDE, 0xD0)
DARK_TEXT   = RGBColor(0x1B, 0x2A, 0x4A)
LIGHT_TEXT  = RGBColor(0xF5, 0xF0, 0xE8)
ACCENT_RED  = RGBColor(0x8B, 0x2F, 0x3A)
CRIMSON     = RGBColor(0xA0, 0x30, 0x40)
GOLD        = RGBColor(0xC4, 0x9A, 0x6C)
BRIGHT_GOLD = RGBColor(0xD4, 0xAA, 0x7C)
GRAY        = RGBColor(0x6B, 0x7B, 0x8D)
DARK_GRAY   = RGBColor(0x4A, 0x5A, 0x6C)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
OLIVE       = RGBColor(0x4A, 0x5D, 0x3A)
TEAL        = RGBColor(0x1A, 0x5C, 0x5A)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
W = prs.slide_width
H = prs.slide_height

# ===== 工具函数 =====
def rect(slide, l, t, w, h, color, alpha=None):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    return s

def oval(slide, l, t, w, h, color):
    s = slide.shapes.add_shape(MSO_SHAPE.OVAL, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    return s

def diamond(slide, l, t, w, h, color):
    s = slide.shapes.add_shape(MSO_SHAPE.DIAMOND, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    return s

def tb(slide, l, t, w, h, text, sz=18, color=DARK_TEXT, bold=False, align=PP_ALIGN.LEFT, font='Microsoft YaHei'):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(sz)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font
    p.alignment = align
    return box

def add_p(tf, text, sz=18, color=DARK_TEXT, bold=False, align=PP_ALIGN.LEFT, font='Microsoft YaHei', spc=Pt(6)):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(sz)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font
    p.alignment = align
    p.space_before = spc
    return p

def bg_navy(slide):
    rect(slide, Inches(0), Inches(0), W, H, DEEP_NAVY)

def bg_cream(slide):
    rect(slide, Inches(0), Inches(0), W, H, PARCHMENT)

def top_bar(slide, color=GOLD):
    rect(slide, Inches(0), Inches(0), W, Inches(0.06), color)

def bottom_bar(slide, color=GOLD):
    rect(slide, Inches(0), H - Inches(0.06), W, Inches(0.06), color)

def header(slide, chapter, title, bg='navy'):
    if bg == 'navy':
        bg_navy(slide)
        top_bar(slide)
        rect(slide, Inches(0), Inches(0), W, Inches(1.3), NAVY)
        tb(slide, Inches(0.8), Inches(0.15), Inches(3), Inches(0.4), chapter, 15, GOLD, True)
        tb(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.7), title, 34, CREAM, True)
        rect(slide, Inches(0.8), Inches(1.25), Inches(4), Inches(0.04), GOLD)
    else:
        bg_cream(slide)
        rect(slide, Inches(0), Inches(0), W, Inches(1.3), NAVY)
        tb(slide, Inches(0.8), Inches(0.15), Inches(3), Inches(0.4), chapter, 15, GOLD, True)
        tb(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.7), title, 34, CREAM, True)
        rect(slide, Inches(0.8), Inches(1.25), Inches(4), Inches(0.04), GOLD)

def card(slide, l, t, w, h, bg_color, border_color=None):
    s = rect(slide, l, t, w, h, bg_color)
    if border_color:
        s.line.color.rgb = border_color
        s.line.width = Pt(1)
    return s

# ====================================================================
# SLIDE 1: 封面
# ====================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
bg_navy(slide=s)

# 装饰：古纹边框效果
rect(s, Inches(0.4), Inches(0.4), Inches(12.5), Inches(6.7), DARK_NAVY)
rect(s, Inches(0.5), Inches(0.5), Inches(12.3), Inches(6.5), DEEP_NAVY)

# 左侧竖条装饰
rect(s, Inches(1.2), Inches(1.2), Inches(0.08), Inches(5.1), GOLD)
rect(s, Inches(1.35), Inches(1.4), Inches(0.03), Inches(4.7), BRIGHT_GOLD)

# 右下角装饰菱形
diamond(s, Inches(10.5), Inches(5.5), Inches(1.2), Inches(1.2), RGBColor(0x25, 0x38, 0x58))
diamond(s, Inches(10.7), Inches(5.7), Inches(0.8), Inches(0.8), RGBColor(0x2F, 0x45, 0x68))

# 标题
tb(s, Inches(1.8), Inches(1.6), Inches(9), Inches(1.0), '天朝的崩溃', 58, CREAM, True)
tb(s, Inches(1.8), Inches(2.8), Inches(9), Inches(0.7), '鸦片战争再研究', 38, GOLD, False)
rect(s, Inches(1.8), Inches(3.6), Inches(6), Inches(0.03), GOLD)
tb(s, Inches(1.8), Inches(3.9), Inches(9), Inches(0.5), '茅海建  著', 22, WARM_CREAM)
tb(s, Inches(1.8), Inches(4.6), Inches(9), Inches(0.5), '生活·读书·新知三联书店  ·  三联·哈佛燕京学术丛书', 14, GRAY)

# 右侧引用
card(s, Inches(8.0), Inches(4.5), Inches(4.5), Inches(2.2), RGBColor(0x15, 0x22, 0x3D))
box = tb(s, Inches(8.3), Inches(4.7), Inches(4.0), Inches(1.8),
         '"鸦片战争是中国历史的转折，\n提出了中国必须近代化的历史使命。\n中国的现代化一日未完成，\n鸦片战争的意义就一分不会减。"', 13, GOLD, False, PP_ALIGN.LEFT)
box.text_frame.paragraphs[0].font.italic = True

# ====================================================================
# SLIDE 2: 作者简介
# ====================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
header(s, '', '作者与写作背景', 'cream')

# 左侧：作者信息
card(s, Inches(0.5), Inches(1.6), Inches(6), Inches(5.4), NAVY)
box = tb(s, Inches(0.8), Inches(1.8), Inches(5.5), Inches(5.0), '', 16, CREAM)
tf = box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = '茅海建'
p.font.size = Pt(30)
p.font.color.rgb = GOLD
p.font.bold = True
p.font.name = 'Microsoft YaHei'

add_p(tf, '', 8, CREAM)
add_p(tf, '华东师范大学历史系教授', 18, CREAM, True)
add_p(tf, '北京大学历史学系兼职教授', 18, CREAM)
add_p(tf, '', 8, CREAM)
add_p(tf, '学术履历', 20, GOLD, True)
add_p(tf, '中山大学历史系 → 华东师范大学（硕士，师从陈旭麓）', 14, CREAM)
add_p(tf, '军事科学院助理研究员', 14, CREAM)
add_p(tf, '中国社会科学院近代史研究所研究员', 14, CREAM)
add_p(tf, '北京大学历史学系教授', 14, CREAM)
add_p(tf, '', 8, CREAM)
add_p(tf, '主要著作', 20, GOLD, True)
add_p(tf, '《天朝的崩溃》1995  ·  《苦命天子》1995', 14, CREAM)
add_p(tf, '《近代的尺度》1998  ·  《戊戌变法史事考》2005', 14, CREAM)
add_p(tf, '《从甲午到戊戌》2009  ·  《依然如旧的月色》2014', 14, CREAM)

# 右侧：写作背景
card(s, Inches(6.8), Inches(1.6), Inches(6), Inches(5.4), WARM_CREAM)
box = tb(s, Inches(7.1), Inches(1.8), Inches(5.5), Inches(5.0), '', 16, DARK_TEXT)
tf = box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = '写作缘起'
p.font.size = Pt(22)
p.font.color.rgb = NAVY
p.font.bold = True
p.font.name = 'Microsoft YaHei'

add_p(tf, '', 8, DARK_TEXT)
add_p(tf, '献给导师陈旭麓教授', 18, ACCENT_RED, True)
add_p(tf, '导师指导的第一篇论文：《鸦片战争时期中英兵力》', 15, DARK_TEXT)
add_p(tf, '1988年导师仙逝，无缘索序', 15, DARK_TEXT)
add_p(tf, '', 8, DARK_TEXT)
add_p(tf, '搜集史料超10年，展纸动笔2年', 16, NAVY, True)
add_p(tf, '1992年初推开一切，整整两年', 15, DARK_TEXT)
add_p(tf, '尝到了著书人都经受过的酸苦辣（没有感到甜）', 15, DARK_TEXT)
add_p(tf, '', 8, DARK_TEXT)
add_p(tf, '方法论', 20, NAVY, True)
add_p(tf, '• 注重人物命运——将纪传体优长融入章节体', 15, DARK_TEXT)
add_p(tf, '• 解释历史现象——用当时人的观念解释当事人的思想和行为', 15, DARK_TEXT)
add_p(tf, '• 历史研究排斥感情的羼入，强调冷静和客观', 15, DARK_TEXT)
add_p(tf, '• 求真毕竟是治史者不灭的梦境', 15, ACCENT_RED, True)

# ====================================================================
# SLIDE 3: 目录
# ====================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
bg_navy(s)
top_bar(s)

tb(s, Inches(0.8), Inches(0.4), Inches(10), Inches(0.8), '全书脉络', 40, CREAM, True)
rect(s, Inches(0.8), Inches(1.2), Inches(3), Inches(0.04), GOLD)

toc = [
    ('绪论', '由琦善卖国而想到的', '打破善恶忠奸的二元框架'),
    ('第一章', '清朝的军事力量', '武器装备·兵力对比·制度缺陷'),
    ('第二章', '骤然而至的战争', '禁烟运动·林则徐使粤·定海陷落'),
    ('第三章', '"剿""抚""剿"的回旋', '道光帝的决策摇摆·初战溃败'),
    ('第四章', '广州的"战局"', '杨芳的荒诞·奕山的谎言'),
    ('第五章', '东南壁垒的倾塌', '颜伯焘与裕谦·璞鼎查北上'),
    ('第六章', '"抚"议再起', '刘韵珂的"十可虑"·战与和的困境'),
    ('第七章', '平等与不平等', '南京条约·天朝的价值标准'),
    ('第八章', '历史的诉说', '中日比较·失败后的奋发'),
]

y = 1.6
for i, (ch, title, desc) in enumerate(toc):
    bg = RGBColor(0x15, 0x22, 0x3D) if i % 2 == 0 else RGBColor(0x1A, 0x2D, 0x50)
    rect(s, Inches(0.8), Inches(y), Inches(11.7), Inches(0.55), bg)
    # 章节号
    rect(s, Inches(0.8), Inches(y), Inches(0.06), Inches(0.55), GOLD)
    tb(s, Inches(1.1), Inches(y + 0.05), Inches(2.2), Inches(0.45), ch, 14, GOLD, True)
    tb(s, Inches(3.5), Inches(y + 0.05), Inches(4.5), Inches(0.45), title, 16, CREAM)
    tb(s, Inches(8.2), Inches(y + 0.05), Inches(4.0), Inches(0.45), desc, 13, GRAY)
    y += 0.6

# ====================================================================
# SLIDE 4: 绪论 - 琦善其人
# ====================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
header(s, '绪论', '由琦善卖国而想到的')

# 左卡片：琦善其人
card(s, Inches(0.5), Inches(1.6), Inches(6), Inches(5.4), RGBColor(0x15, 0x22, 0x3D))
box = tb(s, Inches(0.8), Inches(1.8), Inches(5.5), Inches(5.0), '', 16, CREAM)
tf = box.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]
p.text = '琦善其人'; p.font.size = Pt(22); p.font.color.rgb = GOLD; p.font.bold = True; p.font.name = 'Microsoft YaHei'

add_p(tf, '', 6, CREAM)
add_p(tf, '满洲贵族出身，祖上封一等侯爵', 15, CREAM)
add_p(tf, '16岁荫生入仕，29岁任河南巡抚', 15, CREAM)
add_p(tf, '历山东巡抚、两江总督、成都将军', 15, CREAM)
add_p(tf, '1831年迁直隶总督，1838年擢文渊阁大学士', 15, CREAM)
add_p(tf, '', 6, CREAM)
add_p(tf, '"位极人臣，圣眷正隆"', 18, GOLD, True)
add_p(tf, '', 8, CREAM)
add_p(tf, '为官特点', 18, GOLD, True)
add_p(tf, '• 好用诡道怪行，但也多验明效', 15, CREAM)
add_p(tf, '• 为人傲慢气盛，但官场结交甚广', 15, CREAM)
add_p(tf, '• 勇于任事，好大喜功', 15, CREAM)
add_p(tf, '• 道光帝看重他敢于闯创、敢于负责', 15, CREAM)

# 右卡片：卖国之问
card(s, Inches(6.8), Inches(1.6), Inches(6), Inches(5.4), RGBColor(0x2A, 0x15, 0x1A))
box = tb(s, Inches(7.1), Inches(1.8), Inches(5.5), Inches(5.0), '', 16, CREAM)
tf = box.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]
p.text = '琦善果真卖国吗？'; p.font.size = Pt(22); p.font.color.rgb = CRIMSON; p.font.bold = True; p.font.name = 'Microsoft YaHei'

add_p(tf, '', 6, CREAM)
add_p(tf, '心理动机之疑', 18, GOLD, True)
add_p(tf, '• 世受国恩，无理由背叛', 15, CREAM)
add_p(tf, '• 道光帝待其不薄，没有背叛理由', 15, CREAM)
add_p(tf, '• 与汪精卫因政治不得意而改换门庭不同', 15, CREAM)
add_p(tf, '', 6, CREAM)
add_p(tf, '"贿和"之疑', 18, GOLD, True)
add_p(tf, '• 家赀丰裕，不会见洋货便心旌荡漾', 15, CREAM)
add_p(tf, '• 道光帝亲自审讯，未获证据', 15, CREAM)
add_p(tf, '• 英方文件否认行贿', 15, CREAM)
add_p(tf, '• "世上又哪有强盗上门先行贿后动手的事情"', 15, CREAM)
add_p(tf, '', 6, CREAM)
add_p(tf, '道德审判替代了事实分析', 16, CRIMSON, True)

# ====================================================================
# SLIDE 5: 绪论 - 天朝世界观
# ====================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
bg_cream(s)
rect(s, Inches(0), Inches(0), W, Inches(1.3), NAVY)
tb(s, Inches(0.8), Inches(0.15), Inches(3), Inches(0.4), '绪论', 15, GOLD, True)
tb(s, Inches(0.8), Inches(0.5), Inches(11), Inches(0.7), '善恶忠奸之外：重新审视历史人物', 34, CREAM, True)
rect(s, Inches(0.8), Inches(1.25), Inches(4), Inches(0.04), GOLD)

# 核心观点卡片
card(s, Inches(0.5), Inches(1.6), Inches(5.8), Inches(2.5), NAVY)
box = tb(s, Inches(0.8), Inches(1.8), Inches(5.3), Inches(2.0), '', 16, CREAM)
tf = box.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]
p.text = '传统史学的局限'; p.font.size = Pt(20); p.font.color.rgb = GOLD; p.font.bold = True; p.font.name = 'Microsoft YaHei'

add_p(tf, '', 6, CREAM)
add_p(tf, '中国的历史学，最注重人物评价。打开史籍，善恶忠奸分明，好人坏人一目了然。', 15, CREAM)
add_p(tf, '', 6, CREAM)
add_p(tf, '问题：用道德判断替代事实分析，用"卖国"标签遮蔽历史复杂性。', 15, CRIMSON, True)

card(s, Inches(6.8), Inches(1.6), Inches(6), Inches(2.5), WARM_CREAM)
box = tb(s, Inches(7.1), Inches(1.8), Inches(5.5), Inches(2.0), '', 16, DARK_TEXT)
tf = box.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]
p.text = '茅海建的方法'; p.font.size = Pt(20); p.font.color.rgb = NAVY; p.font.bold = True; p.font.name = 'Microsoft YaHei'

add_p(tf, '', 6, DARK_TEXT)
add_p(tf, '• 回到历史现场，理解当事人的处境与选择', 15, DARK_TEXT)
add_p(tf, '• 用当时人的观念解释当时人的行为', 15, DARK_TEXT)
add_p(tf, '• 超越简单的道德审判', 15, DARK_TEXT)
add_p(tf, '• 在"天朝"的世界观中理解"天朝"的人', 15, DARK_TEXT)

# 底部：天朝世界观
card(s, Inches(0.5), Inches(4.4), Inches(12.3), Inches(2.8), PARCHMENT)
box = tb(s, Inches(0.8), Inches(4.6), Inches(11.8), Inches(2.4), '', 16, DARK_TEXT)
tf = box.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]
p.text = '19世纪的"天朝"：一个特殊的"世界"'; p.font.size = Pt(22); p.font.color.rgb = NAVY; p.font.bold = True; p.font.name = 'Microsoft YaHei'

add_p(tf, '', 6, DARK_TEXT)
add_p(tf, '鸦片战争之前，中华文明一直是相对独立地发展的，并以其优越性向外输出，在东亚地区形成了以中国为中心的汉文化圈。', 15, DARK_TEXT)
add_p(tf, '', 4, DARK_TEXT)
add_p(tf, '生活在"天朝"中的人们，自有一套迥然相别的价值标准，另有一种平等观念。他们对今天看来为"平等"的条款往往愤愤不平，而对今天看来为"不平等"的待遇却浑然不觉。', 15, DARK_TEXT)
add_p(tf, '', 4, DARK_TEXT)
add_p(tf, '中华文明不是不想了解西方，而是自认为已经了解了西方——"天下"就这么大，"天朝"就在中央，其余不过是化外之邦。', 15, ACCENT_RED, True)

# ====================================================================
# SLIDE 6: 第一章 - 武器装备（清军）
# ====================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
header(s, '第一章', '清朝的军事力量：清军武器装备')

# 大卡片
card(s, Inches(0.5), Inches(1.6), Inches(12.3), Inches(5.4), RGBColor(0x15, 0x22, 0x3D))

# 左：火器
box = tb(s, Inches(0.8), Inches(1.8), Inches(5.5), Inches(5.0), '', 16, CREAM)
tf = box.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]
p.text = '火器：自制的老式"洋枪洋炮"'; p.font.size = Pt(20); p.font.color.rgb = GOLD; p.font.bold = True; p.font.name = 'Microsoft YaHei'

add_p(tf, '', 6, CREAM)
add_p(tf, '鸟枪（主要单兵火器）', 18, CREAM, True)
add_p(tf, '• 仿1548年葡萄牙火绳枪，落后英军200余年', 15, CREAM)
add_p(tf, '• 前装滑膛火绳枪，枪长2.01米', 15, CREAM)
add_p(tf, '• 射程约100米，射速1-2发/分钟', 15, CREAM)
add_p(tf, '• 种类达58种，大同小异', 15, CREAM)
add_p(tf, '• 无定期更换制度，有的使用166年', 15, CREAM)
add_p(tf, '', 6, CREAM)
add_p(tf, '火炮', 18, CREAM, True)
add_p(tf, '• 仿17-18世纪西方加农炮', 15, CREAM)
add_p(tf, '• 铁质差，气孔多，易炸裂', 15, CREAM)
add_p(tf, '• 泥模工艺，炮膛粗糙', 15, CREAM)
add_p(tf, '• 无瞄准器具，靠经验射击', 15, CREAM)
add_p(tf, '• 炮弹仅实心弹一种', 15, CREAM)

# 右：冷兵器与制度
box = tb(s, Inches(7.0), Inches(1.8), Inches(5.5), Inches(5.0), '', 16, CREAM)
tf = box.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]
p.text = '冷兵器与制度缺陷'; p.font.size = Pt(20); p.font.color.rgb = GOLD; p.font.bold = True; p.font.name = 'Microsoft YaHei'

add_p(tf, '', 6, CREAM)
add_p(tf, '冷热混用', 18, CREAM, True)
add_p(tf, '• 鸟枪手与刀矛弓箭手比例约5:5', 15, CREAM)
add_p(tf, '• 鸟枪太长无法装枪刺', 15, CREAM)
add_p(tf, '• 短兵相接难以应敌', 15, CREAM)
add_p(tf, '', 6, CREAM)
add_p(tf, '制度性衰败', 18, CREAM, True)
add_p(tf, '• 承平日久，军费不足', 15, CREAM)
add_p(tf, '• 无定期修造报废更换制度', 15, CREAM)
add_p(tf, '• 赶制火器质量尤其低劣', 15, CREAM)
add_p(tf, '', 6, CREAM)
add_p(tf, '"清军使用的是自制的老式的\'洋枪洋炮\'。就型制样式而言，与英军相比，整整落后了二百余年。"', 14, GOLD, False, PP_ALIGN.LEFT)
add_p(tf, '', 4, CREAM)
add_p(tf, '火药和管型火器都是中国发明的，但中国一直处于前科学时期，没有形成科学理论和实验体系，使得中国火器的发展受到了根本性的制约。', 14, GRAY)

# ====================================================================
# SLIDE 7: 第一章 - 英军装备与对比
# ====================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
header(s, '第一章', '中英军事力量对比', 'cream')

# 对比表格
headers = ['对比项', '清军', '英军', '差距']
rows = [
    ['时代', '冷热兵器混用', '初步火器时代', '200+年'],
    ['步枪射程', '~100米', '200-300米', '2-3倍'],
    ['步枪射速', '1-2发/分', '3-4发/分', '2-3倍'],
    ['火炮工艺', '泥模/铁质差', '铁模/镗床加工', '代差'],
    ['炮弹种类', '仅实心弹', '实心/霰弹/爆破弹', '3种'],
    ['瞄准器具', '无或仅有星斗', '完善的瞄准系统', '质差'],
    ['枪刺', '无（枪身太长）', '标配', '缺失'],
    ['换装制度', '无定期', '定期换装', '制度差'],
]

y = 1.5
# Header row
rect(s, Inches(0.5), Inches(y), Inches(3), Inches(0.5), NAVY)
tb(s, Inches(0.5), Inches(y+0.05), Inches(3), Inches(0.4), '对比项', 14, CREAM, True, PP_ALIGN.CENTER)
rect(s, Inches(3.6), Inches(y), Inches(3.2), Inches(0.5), ACCENT_RED)
tb(s, Inches(3.6), Inches(y+0.05), Inches(3.2), Inches(0.4), '清军', 14, WHITE, True, PP_ALIGN.CENTER)
rect(s, Inches(6.9), Inches(y), Inches(3.2), Inches(0.5), OLIVE)
tb(s, Inches(6.9), Inches(y+0.05), Inches(3.2), Inches(0.4), '英军', 14, WHITE, True, PP_ALIGN.CENTER)
rect(s, Inches(10.2), Inches(y), Inches(2.6), Inches(0.5), DARK_GRAY)
tb(s, Inches(10.2), Inches(y+0.05), Inches(2.6), Inches(0.4), '差距', 14, CREAM, True, PP_ALIGN.CENTER)

y += 0.55
for i, row in enumerate(rows):
    bg = WARM_CREAM if i % 2 == 0 else CREAM
    rect(s, Inches(0.5), Inches(y), Inches(3), Inches(0.48), bg)
    tb(s, Inches(0.5), Inches(y+0.04), Inches(3), Inches(0.4), row[0], 13, NAVY, True, PP_ALIGN.CENTER)
    rect(s, Inches(3.6), Inches(y), Inches(3.2), Inches(0.48), bg)
    tb(s, Inches(3.6), Inches(y+0.04), Inches(3.2), Inches(0.4), row[1], 13, ACCENT_RED, False, PP_ALIGN.CENTER)
    rect(s, Inches(6.9), Inches(y), Inches(3.2), Inches(0.48), bg)
    tb(s, Inches(6.9), Inches(y+0.04), Inches(3.2), Inches(0.4), row[2], 13, OLIVE, False, PP_ALIGN.CENTER)
    rect(s, Inches(10.2), Inches(y), Inches(2.6), Inches(0.48), bg)
    tb(s, Inches(10.2), Inches(y+0.04), Inches(2.6), Inches(0.4), row[3], 13, CRIMSON, True, PP_ALIGN.CENTER)
    y += 0.5

# 结论
card(s, Inches(0.5), Inches(5.8), Inches(12.3), Inches(1.2), NAVY)
box = tb(s, Inches(0.8), Inches(5.9), Inches(11.8), Inches(1.0), '', 16, CREAM)
tf = box.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]
p.text = '核心结论'; p.font.size = Pt(18); p.font.color.rgb = GOLD; p.font.bold = True; p.font.name = 'Microsoft YaHei'
add_p(tf, '如果把这些枪改换成持枪的士兵，多少名清军士兵方能抵得上一名英军士兵？——茅海建', 15, CREAM)
add_p(tf, '以血肉之躯对抗工业文明，败局早已注定。', 15, CRIMSON, True)

# ====================================================================
# SLIDE 8: 第二章 - 禁烟运动
# ====================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
header(s, '第二章', '骤然而至的战争：禁烟运动')

# 时间线
events = [
    ('1838.06', '黄爵滋上奏', '严禁鸦片，提出"吸烟者诛"'),
    ('1838.10', '29份议复', '19份主张禁烟重点在海口，地方官推卸责任'),
    ('1838.10', '皇室丑闻', '庄亲王、镇国公在尼僧庙吸食鸦片'),
    ('1838.11', '天津大案', '琦善查获鸦片13万两，来源直指广东'),
    ('1839.01', '林则徐南下', '钦差大臣赴粤，"焚香九拜，发传牌，遂起程"'),
    ('1839.06', '虎门销烟', '收缴鸦片237万余斤，历时23天'),
    ('1840.06', '英军抵华', '懿律率舰队到达，战争爆发'),
]

y = 1.6
for i, (date, title, desc) in enumerate(events):
    bg = RGBColor(0x15, 0x22, 0x3D) if i % 2 == 0 else RGBColor(0x1A, 0x2D, 0x50)
    rect(s, Inches(0.5), Inches(y), Inches(12.3), Inches(0.7), bg)
    rect(s, Inches(0.5), Inches(y), Inches(0.06), Inches(0.7), GOLD)
    tb(s, Inches(0.8), Inches(y+0.1), Inches(1.6), Inches(0.5), date, 15, GOLD, True)
    tb(s, Inches(2.6), Inches(y+0.1), Inches(2.5), Inches(0.5), title, 16, CREAM, True)
    tb(s, Inches(5.3), Inches(y+0.1), Inches(7.2), Inches(0.5), desc, 14, GRAY)
    y += 0.72

# 底部分析
card(s, Inches(0.5), Inches(6.7), Inches(12.3), Inches(0.6), NAVY)
tb(s, Inches(0.8), Inches(6.75), Inches(11.8), Inches(0.5),
   '禁烟的重点从沿海扩大到内地，从查禁"夷商"变为全国范围内的捕杀瘾君子的国内司法行动。',
   14, GOLD, False, PP_ALIGN.LEFT)

# ====================================================================
# SLIDE 9: 第二章 - 定海之战
# ====================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
header(s, '第二章', '定海之战：第一场溃败')

# 左：战役经过
card(s, Inches(0.5), Inches(1.6), Inches(6), Inches(5.4), RGBColor(0x15, 0x22, 0x3D))
box = tb(s, Inches(0.8), Inches(1.8), Inches(5.5), Inches(5.0), '', 16, CREAM)
tf = box.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]
p.text = '1840年7月5日 · 定海'; p.font.size = Pt(20); p.font.color.rgb = GOLD; p.font.bold = True; p.font.name = 'Microsoft YaHei'

add_p(tf, '', 6, CREAM)
add_p(tf, '英军舰炮仅用9分钟，基本击毁清军战船和岸炮还击能力。', 15, CREAM)
add_p(tf, '', 6, CREAM)
add_p(tf, '知县姚怀祥登上英舰后说：', 16, GOLD, True)
add_p(tf, '"你们把战争施加于民众身上，而不是我们这些从未伤害过你们的人；我们看到了你们的强大，也知道对抗将是发疯，但我们必须恪尽职守，尽管如此做会遭至失败。"', 14, CREAM)
add_p(tf, '', 6, CREAM)
add_p(tf, '战损对比', 16, GOLD, True)
add_p(tf, '• 清军参战1540人，战死仅13人，受伤13人', 15, CREAM)
add_p(tf, '• 英方宣称：战斗中毫无伤亡', 15, CRIMSON, True)
add_p(tf, '', 6, CREAM)
add_p(tf, '知县姚怀祥投水自尽', 15, CREAM)
add_p(tf, '总兵张朝发中弹落水，后不治', 15, CREAM)

# 右：北京的无知
card(s, Inches(6.8), Inches(1.6), Inches(6), Inches(5.4), RGBColor(0x2A, 0x15, 0x1A))
box = tb(s, Inches(7.1), Inches(1.8), Inches(5.5), Inches(5.0), '', 16, CREAM)
tf = box.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]
p.text = '同一时间的北京'; p.font.size = Pt(20); p.font.color.rgb = GOLD; p.font.bold = True; p.font.name = 'Microsoft YaHei'

add_p(tf, '', 8, CREAM)
add_p(tf, '道光帝：', 18, CREAM, True)
add_p(tf, '例行"诣绮春园问皇太后安"，然后回銮处理日常公文。', 15, CREAM)
add_p(tf, '', 12, CREAM)
add_p(tf, '曾国藩：', 18, CREAM, True)
add_p(tf, '因客来访耽误了读书，在日记中狠狠自责，自励须"日日用功有常"，以能"文章报国"。', 15, CREAM)
add_p(tf, '', 12, CREAM)
add_p(tf, '"正当道光帝享以清静时，正当后来以武功名扬天下的曾国藩琢磨\'文章报国\'之道时，远去北京数千里的浙江省定海县，已是一片炮声隆隆，笼罩于呛人的硝烟之中了。"', 14, GOLD)
add_p(tf, '', 8, CREAM)
add_p(tf, '信息的鸿沟，决定了决策的盲目。', 16, CRIMSON, True)

# ====================================================================
# SLIDE 10: 第三章 - "剿""抚"回旋
# ====================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
header(s, '第三章', '"剿""抚""剿"的回旋')

# 三段式流程
labels = ['剿', '抚', '剿']
descs = [
    '初战：定海陷落\n1840年7月英军攻占舟山\n道光帝震怒，决定"剿夷"',
    '琦善谈判：天津交涉\n英军北上大沽口\n琦善"抚夷"交涉失败',
    '再战：虎门陷落\n琦善革职锁京\n杨芳、奕山接替',
]
colors = [ACCENT_RED, GOLD, ACCENT_RED]

x = 0.8
for i in range(3):
    # 大卡片
    card(s, Inches(x), Inches(1.6), Inches(3.5), Inches(2.8), colors[i])
    tb(s, Inches(x+0.3), Inches(1.7), Inches(2.9), Inches(0.7), labels[i], 40, WHITE, True, PP_ALIGN.CENTER)
    tb(s, Inches(x+0.3), Inches(2.5), Inches(2.9), Inches(1.5), descs[i], 13, WHITE, False, PP_ALIGN.CENTER)
    # 箭头
    if i < 2:
        tb(s, Inches(x+3.6), Inches(2.5), Inches(0.6), Inches(0.5), '→', 28, GRAY, True, PP_ALIGN.CENTER)
    x += 4.2

# 底部：核心洞察
card(s, Inches(0.5), Inches(4.7), Inches(12.3), Inches(2.5), WARM_CREAM)
box = tb(s, Inches(0.8), Inches(4.9), Inches(11.8), Inches(2.0), '', 16, DARK_TEXT)
tf = box.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]
p.text = '核心洞察'; p.font.size = Pt(22); p.font.color.rgb = NAVY; p.font.bold = True; p.font.name = 'Microsoft YaHei'

add_p(tf, '', 6, DARK_TEXT)
add_p(tf, '• "剿捕档"——军机处将鸦片战争等同于平定叛乱，反映了"天朝"的天下观念', 16, DARK_TEXT)
add_p(tf, '• 传统的御外攘夷武库中，只有"剿"和"抚"两套程序', 16, DARK_TEXT)
add_p(tf, '• 道光帝交并轮番操之上阵，一波三折，分寸大乱', 16, DARK_TEXT)
add_p(tf, '• "决策者自然有权多变，但每一变都会在战场上付出相应的代价。"', 16, ACCENT_RED, True)

# ====================================================================
# SLIDE 11: 第四章 - 广州战局
# ====================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
header(s, '第四章', '广州的"战局"：一个骗局')

# 钳口令
card(s, Inches(0.5), Inches(1.6), Inches(12.3), Inches(1.8), RGBColor(0x2A, 0x15, 0x1A))
box = tb(s, Inches(0.8), Inches(1.7), Inches(11.8), Inches(1.5), '', 16, CREAM)
tf = box.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]
p.text = '道光帝的"钳口令"'; p.font.size = Pt(20); p.font.color.rgb = GOLD; p.font.bold = True; p.font.name = 'Microsoft YaHei'

add_p(tf, '', 4, CREAM)
add_p(tf, '（琦善）被人恐吓，奏报粤省情形，妄称地利无要可扼，军械无利可恃，兵力不固，民情不坚。摘举数端，危言要挟，更不知是何肺腑？', 15, CREAM)
add_p(tf, '', 4, CREAM)
add_p(tf, '效果：不仅不许败，而且不许言败。杨芳和奕山面前只有一条出路——捏谎。', 15, CRIMSON, True)

# 杨芳
card(s, Inches(0.5), Inches(3.7), Inches(5.8), Inches(3.5), RGBColor(0x15, 0x22, 0x3D))
box = tb(s, Inches(0.8), Inches(3.9), Inches(5.3), Inches(3.0), '', 16, CREAM)
tf = box.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]
p.text = '杨芳的"果勇"'; p.font.size = Pt(20); p.font.color.rgb = GOLD; p.font.bold = True; p.font.name = 'Microsoft YaHei'

add_p(tf, '', 4, CREAM)
add_p(tf, '• 戎马55载，以平定张格尔封三等果勇侯', 15, CREAM)
add_p(tf, '• 到广州后"终日唯购钟表洋货为事"', 15, CREAM)
add_p(tf, '• "购买马桶御炮"——以邪制邪', 15, CREAM)
add_p(tf, '• 与林则徐14天见面11次', 15, CREAM)
add_p(tf, '• 英军已看见广州城墙', 15, CREAM)
add_p(tf, '', 4, CREAM)
add_p(tf, '杨芳对西方利器的不解，以"马桶""草人""道场""鬼神"来应对。', 14, GRAY)

# 奕山
card(s, Inches(6.8), Inches(3.7), Inches(6), Inches(3.5), RGBColor(0x15, 0x22, 0x3D))
box = tb(s, Inches(7.1), Inches(3.9), Inches(5.5), Inches(3.0), '', 16, CREAM)
tf = box.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]
p.text = '奕山的谎言'; p.font.size = Pt(20); p.font.color.rgb = GOLD; p.font.bold = True; p.font.name = 'Microsoft YaHei'

add_p(tf, '', 4, CREAM)
add_p(tf, '• 捏谎报捷，欺骗道光帝', 15, CREAM)
add_p(tf, '• 广州到北京的河川山岭，成为谎话的天然屏障', 15, CREAM)
add_p(tf, '• 整个广州战局，完全成为一个骗局', 15, CREAM)
add_p(tf, '', 4, CREAM)
add_p(tf, '璞鼎查到达后，奕山再次行骗：', 16, GOLD, True)
add_p(tf, '• 璞鼎查要求与清方"全权"大臣谈判', 15, CREAM)
add_p(tf, '• 奕山等人明白无误地知道这一切', 15, CREAM)
add_p(tf, '• 但为了遮盖先前的谎言，选择隐瞒', 15, CRIMSON, True)

# ====================================================================
# SLIDE 12: 第五章 - 东南壁垒
# ====================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
header(s, '第五章', '东南壁垒的倾塌')

# 颜伯焘与裕谦
card(s, Inches(0.5), Inches(1.6), Inches(5.8), Inches(2.5), RGBColor(0x15, 0x22, 0x3D))
box = tb(s, Inches(0.8), Inches(1.8), Inches(5.3), Inches(2.0), '', 16, CREAM)
tf = box.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]
p.text = '颜伯焘与裕谦：道光帝心目中的长城'; p.font.size = Pt(18); p.font.color.rgb = GOLD; p.font.bold = True; p.font.name = 'Microsoft YaHei'

add_p(tf, '', 4, CREAM)
add_p(tf, '• 抗战言论最坚决，深孚清望', 15, CREAM)
add_p(tf, '• 筹防措施最彻底，建起坚固壁垒', 15, CREAM)
add_p(tf, '• 颜伯焘坐镇厦门，裕谦长驻镇海', 15, CREAM)
add_p(tf, '', 4, CREAM)
add_p(tf, '结果：当英国军舰鼓浪而来时，东南的壁垒倾塌了。', 15, CRIMSON, True)

# 璞鼎查
card(s, Inches(6.8), Inches(1.6), Inches(6), Inches(2.5), RGBColor(0x1A, 0x2D, 0x50))
box = tb(s, Inches(7.1), Inches(1.8), Inches(5.5), Inches(2.0), '', 16, CREAM)
tf = box.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]
p.text = '璞鼎查的东来'; p.font.size = Pt(18); p.font.color.rgb = GOLD; p.font.bold = True; p.font.name = 'Microsoft YaHei'

add_p(tf, '', 4, CREAM)
add_p(tf, '• 从伦敦到澳门仅用67天，破纪录', 15, CREAM)
add_p(tf, '• 比林则徐(61天)、琦善(56天)、奕山(57天)还快', 15, CREAM)
add_p(tf, '• 科学缩短了空间的距离', 15, CREAM)
add_p(tf, '• 与义律迥然不同的强硬风格', 15, CREAM)

# 底部：英军北上路线
card(s, Inches(0.5), Inches(4.4), Inches(12.3), Inches(2.8), WARM_CREAM)
box = tb(s, Inches(0.8), Inches(4.6), Inches(11.8), Inches(2.4), '', 16, DARK_TEXT)
tf = box.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]
p.text = '英军北上：连陷东南'; p.font.size = Pt(20); p.font.color.rgb = NAVY; p.font.bold = True; p.font.name = 'Microsoft YaHei'

add_p(tf, '', 4, DARK_TEXT)
add_p(tf, '厦门（1841.8）→ 定海（1841.10）→ 镇海（1841.10）→ 宁波（1841.10）→ 吴淞（1842.6）→ 镇江（1842.7）→ 南京（1842.8）', 16, ACCENT_RED, True)
add_p(tf, '', 4, DARK_TEXT)
add_p(tf, '• 三总兵战死，裕谦自杀', 15, DARK_TEXT)
add_p(tf, '• 颜伯焘苦心经营的厦门壁垒，一日即破', 15, DARK_TEXT)
add_p(tf, '• 东南沿海千里防线，全面崩溃', 15, DARK_TEXT)
add_p(tf, '', 4, DARK_TEXT)
add_p(tf, '"150多年来，将失败归结于好色、贪货、抗敌意志不坚定等道德上的非难，使得人们长久地未究诘事理，幻想着制\'夷\'的英雄。战争失败的必然性，并没有因为各停战协定而明朗。"', 14, NAVY)

# ====================================================================
# SLIDE 13: 第六章 - 刘韵珂
# ====================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
header(s, '第六章', '"抚"议再起：刘韵珂与"十可虑"')

# 人物卡片
card(s, Inches(0.5), Inches(1.6), Inches(5.8), Inches(2.5), RGBColor(0x15, 0x22, 0x3D))
box = tb(s, Inches(0.8), Inches(1.8), Inches(5.3), Inches(2.0), '', 16, CREAM)
tf = box.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]
p.text = '刘韵珂其人'; p.font.size = Pt(20); p.font.color.rgb = GOLD; p.font.bold = True; p.font.name = 'Microsoft YaHei'

add_p(tf, '', 4, CREAM)
add_p(tf, '• 不是翰林、进士，仅拔贡生出身', 15, CREAM)
add_p(tf, '• 14年间由七品小京官升至浙江巡抚', 15, CREAM)
add_p(tf, '• 办事结实 + 为人乖巧', 15, CREAM)
add_p(tf, '• 在讲究学历、门第的道光朝，可视作特例', 15, CREAM)

# 转变
card(s, Inches(6.8), Inches(1.6), Inches(6), Inches(2.5), RGBColor(0x2A, 0x15, 0x1A))
box = tb(s, Inches(7.1), Inches(1.8), Inches(5.5), Inches(2.0), '', 16, CREAM)
tf = box.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]
p.text = '从主"剿"到主"抚"'; p.font.size = Pt(20); p.font.color.rgb = GOLD; p.font.bold = True; p.font.name = 'Microsoft YaHei'

add_p(tf, '', 4, CREAM)
add_p(tf, '• 战争初：坚定的主战派，与林则徐朝夕相处', 15, CREAM)
add_p(tf, '• 定海、镇海陷落后：惊骇失色', 15, CREAM)
add_p(tf, '• "战、守、抚三端，今战、守不利，抚又不可"', 15, CRIMSON, True)
add_p(tf, '• 讲真话，需要点勇气，也需要点正气', 15, CREAM)

# 十可虑
card(s, Inches(0.5), Inches(4.4), Inches(12.3), Inches(2.8), WARM_CREAM)
box = tb(s, Inches(0.8), Inches(4.5), Inches(5.5), Inches(2.5), '', 16, DARK_TEXT)
tf = box.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]
p.text = '"十可虑"精选'; p.font.size = Pt(20); p.font.color.rgb = NAVY; p.font.bold = True; p.font.name = 'Microsoft YaHei'

add_p(tf, '', 4, DARK_TEXT)
add_p(tf, '一、对英条约签订后，其他国家望而效尤，怎么办？', 14, DARK_TEXT)
add_p(tf, '三、英国屡言北上天津，如何"能杜其北上之心"？', 14, DARK_TEXT)
add_p(tf, '五、民夷争讼，英方拒不交凶，怎么办？', 14, DARK_TEXT)
add_p(tf, '七、赦免"汉奸"后，匪徒投靠英方扰民，怎么办？', 14, DARK_TEXT)
add_p(tf, '九、英人"建造夷楼""大有据邑之意"，怎么办？', 14, DARK_TEXT)

box2 = tb(s, Inches(7), Inches(4.5), Inches(5.5), Inches(2.5), '', 16, DARK_TEXT)
tf2 = box2.text_frame; tf2.word_wrap = True
p2 = tf2.paragraphs[0]
p2.text = '意义'; p2.font.size = Pt(20); p2.font.color.rgb = NAVY; p2.font.bold = True; p2.font.name = 'Microsoft YaHei'

add_p(tf2, '', 4, DARK_TEXT)
add_p(tf2, '这是一篇新的"十可虑"，是对战后中外关系的深层次思考。', 15, DARK_TEXT)
add_p(tf2, '', 4, DARK_TEXT)
add_p(tf2, '本无国际知识的刘韵珂，所提出的问题以今日之眼光观之十分可笑，但其中蕴含的忧虑，许多在后来不幸言中。', 15, DARK_TEXT)
add_p(tf2, '', 4, DARK_TEXT)
add_p(tf2, '刘韵珂代表了那个时代最有见识的官员——他们开始意识到"天朝"遇到了前所未有的变局。', 15, ACCENT_RED, True)

# ====================================================================
# SLIDE 14: 第七章 - 南京条约
# ====================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
header(s, '第七章', '平等与不平等：南京条约')

# 蒋廷黻名言
card(s, Inches(0.5), Inches(1.6), Inches(12.3), Inches(1.2), NAVY)
box = tb(s, Inches(0.8), Inches(1.7), Inches(11.8), Inches(1.0), '', 16, CREAM)
tf = box.text_frame; tf.word_wrap = True; tf.paragraphs[0].alignment = PP_ALIGN.CENTER
p = tf.paragraphs[0]
p.text = '"中西关系是特别的。在鸦片战争以前，我们不肯给外国平等待遇；在以后，他们不肯给我们平等待遇。"'
p.font.size = Pt(20); p.font.color.rgb = GOLD; p.font.italic = True; p.font.name = 'Microsoft YaHei'; p.alignment = PP_ALIGN.CENTER
add_p(tf, '——蒋廷黻', 14, GRAY, False, PP_ALIGN.CENTER)

# 不平等条款
card(s, Inches(0.5), Inches(3.1), Inches(5.8), Inches(2.0), RGBColor(0x2A, 0x15, 0x1A))
box = tb(s, Inches(0.8), Inches(3.2), Inches(5.3), Inches(1.8), '', 16, CREAM)
tf = box.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]
p.text = '明确的不平等条款'; p.font.size = Pt(18); p.font.color.rgb = GOLD; p.font.bold = True; p.font.name = 'Microsoft YaHei'

add_p(tf, '', 4, CREAM)
add_p(tf, '• 割地（香港）', 15, CREAM)
add_p(tf, '• 赔款（鸦片、商欠、军费）', 15, CREAM)
add_p(tf, '• 赦免"汉奸"', 15, CREAM)

# 争议性条款
card(s, Inches(6.8), Inches(3.1), Inches(6), Inches(2.0), RGBColor(0x1A, 0x2D, 0x50))
box = tb(s, Inches(7.1), Inches(3.2), Inches(5.5), Inches(1.8), '', 16, CREAM)
tf = box.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]
p.text = '争议性条款'; p.font.size = Pt(18); p.font.color.rgb = GOLD; p.font.bold = True; p.font.name = 'Microsoft YaHei'

add_p(tf, '', 4, CREAM)
add_p(tf, '• 五口通商——打破一口通商的束缚', 15, CREAM)
add_p(tf, '• 废除行商——消除灰色交易', 15, CREAM)
add_p(tf, '• 新定税则——从短期负面到长期正面', 15, CREAM)

# 当时人的反应
card(s, Inches(0.5), Inches(5.4), Inches(12.3), Inches(1.8), WARM_CREAM)
box = tb(s, Inches(0.8), Inches(5.5), Inches(11.8), Inches(1.5), '', 16, DARK_TEXT)
tf = box.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]
p.text = '当时人的反应'; p.font.size = Pt(18); p.font.color.rgb = NAVY; p.font.bold = True; p.font.name = 'Microsoft YaHei'

add_p(tf, '', 4, DARK_TEXT)
add_p(tf, '李星沅（江苏布政使）："夷妇与大皇帝并书"——本能地感到无法向历史交账。他最看不惯的不是条约内容本身，而是"公然大书特书"。', 14, DARK_TEXT)
add_p(tf, '', 2, DARK_TEXT)
add_p(tf, '关键发现：当时人对"平等"条款愤愤不平，对"不平等"待遇浑然不觉。"天朝"自有一套迥然相别的价值标准。', 14, ACCENT_RED, True)

# ====================================================================
# SLIDE 15: 第八章 - 中日比较
# ====================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
header(s, '第八章', '历史的诉说：中日比较')

# 清朝
card(s, Inches(0.5), Inches(1.6), Inches(5.8), Inches(5.4), NAVY)
box = tb(s, Inches(0.8), Inches(1.8), Inches(5.3), Inches(5.0), '', 16, CREAM)
tf = box.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]
p.text = '清朝：天朝的迷梦'; p.font.size = Pt(22); p.font.color.rgb = GOLD; p.font.bold = True; p.font.name = 'Microsoft YaHei'

add_p(tf, '', 6, CREAM)
add_p(tf, '战后的表现', 18, GOLD, True)
add_p(tf, '• 琦善：定斩监候→释放→复职', 15, CREAM)
add_p(tf, '• 道光帝："何事未曾办过！"', 15, CREAM)
add_p(tf, '• 御史陈庆镛：仍持"人心论"', 15, CREAM)
add_p(tf, '• 琦善最终官复原职', 15, CREAM)
add_p(tf, '', 6, CREAM)
add_p(tf, '"清朝似乎仍未从\'天朝\'的迷梦中醒来，勇敢地进入全新的世界，而是依然如故，就像一切都未发生。"', 14, GOLD)
add_p(tf, '', 6, CREAM)
add_p(tf, '一个失败的民族在战后认真思过，幡然变计，是对殉国者最大的尊崇、最好的纪念。', 14, CRIMSON, True)

# 日本
card(s, Inches(6.8), Inches(1.6), Inches(6), Inches(5.4), OLIVE)
box = tb(s, Inches(7.1), Inches(1.8), Inches(5.5), Inches(5.0), '', 16, CREAM)
tf = box.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]
p.text = '日本：黑船的启示'; p.font.size = Pt(22); p.font.color.rgb = GOLD; p.font.bold = True; p.font.name = 'Microsoft YaHei'

add_p(tf, '', 6, CREAM)
add_p(tf, '1853年 · 培里率"黑船"抵东京湾', 18, CREAM, True)
add_p(tf, '• 震动不亚于英国军舰开抵大沽口', 15, CREAM)
add_p(tf, '• 5年间未抵抗，签订不平等条约', 15, CREAM)
add_p(tf, '• 除割地赔款外，日本"享受"着与中国同等的待遇', 15, CREAM)
add_p(tf, '', 6, CREAM)
add_p(tf, '"名茶上喜选，只消喝四碗，惊破太平梦，彻夜不能眠。"', 16, GOLD, True)
add_p(tf, '', 6, CREAM)
add_p(tf, '战舰化作浓茶，引起神经中枢的高度兴奋，引起日本民族不睡觉的奋斗，引起明治维新。', 15, CREAM)
add_p(tf, '', 6, CREAM)
add_p(tf, '"安政五国条约"的失败是今日日本成功之母。', 16, GOLD, True)
add_p(tf, '失败的民族仍有机会再度辉煌，关键在于战后的奋发。', 15, CREAM)

# ====================================================================
# SLIDE 16: 核心论点
# ====================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
bg_navy(s)
top_bar(s)

tb(s, Inches(0.8), Inches(0.3), Inches(11), Inches(0.7), '本书核心论点', 36, CREAM, True)
rect(s, Inches(0.8), Inches(1.0), Inches(3), Inches(0.04), GOLD)

points = [
    ('武器落后是失败的直接原因', '清军冷热兵器混用，落后英军二百余年，以血肉之躯对抗工业文明。'),
    ('制度腐朽是失败的根本原因', '军费不足、承平日久、无定期换装、军官靠经验瞄准——系统性衰败。'),
    ('信息闭塞导致战略误判', '从皇帝到平民都不知英国力量，沉醉于"天朝"迷梦。'),
    ('道德评判遮蔽历史真相', '琦善"卖国"的标签经不起推敲，需要超越善恶忠奸的二元框架。'),
    ('"剿""抚"困境反映观念局限', '传统武库只有两套程序，无法应对全新的国际格局。'),
    ('战后不反省才是真正的悲剧', '失败不可怕，可怕的是"就像一切都未发生"。'),
]

y = 1.3
for i, (title, desc) in enumerate(points):
    bg = RGBColor(0x15, 0x22, 0x3D) if i % 2 == 0 else RGBColor(0x1A, 0x2D, 0x50)
    rect(s, Inches(0.5), Inches(y), Inches(12.3), Inches(0.85), bg)
    rect(s, Inches(0.5), Inches(y), Inches(0.06), Inches(0.85), GOLD)
    # 序号
    oval(s, Inches(0.8), Inches(y + 0.15), Inches(0.5), Inches(0.5), GOLD)
    tb(s, Inches(0.8), Inches(y + 0.18), Inches(0.5), Inches(0.45), str(i+1), 16, DEEP_NAVY, True, PP_ALIGN.CENTER)
    tb(s, Inches(1.5), Inches(y + 0.05), Inches(10.5), Inches(0.4), title, 18, CREAM, True)
    tb(s, Inches(1.5), Inches(y + 0.45), Inches(10.5), Inches(0.35), desc, 13, GRAY)
    y += 0.92

# ====================================================================
# SLIDE 17: 历史反思
# ====================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
header(s, '', '历史的反思')

# 核心引文
card(s, Inches(0.5), Inches(1.6), Inches(12.3), Inches(2.2), NAVY)
box = tb(s, Inches(0.8), Inches(1.7), Inches(11.8), Inches(2.0), '', 16, CREAM)
tf = box.text_frame; tf.word_wrap = True; tf.paragraphs[0].alignment = PP_ALIGN.CENTER
p = tf.paragraphs[0]
p.text = '"对于列强的入侵，武力抵抗无疑是正确的；\n但这种抵抗注定要失败，另作选择也是明智的。\n前者是道德层面的，后者是政治层面的。"'
p.font.size = Pt(24); p.font.color.rgb = GOLD; p.font.italic = True; p.font.name = 'Microsoft YaHei'; p.alignment = PP_ALIGN.CENTER
add_p(tf, '', 8, CREAM, False, PP_ALIGN.CENTER)
add_p(tf, '——茅海建', 16, GRAY, False, PP_ALIGN.CENTER)

# 启示
card(s, Inches(0.5), Inches(4.1), Inches(12.3), Inches(3.0), WARM_CREAM)
box = tb(s, Inches(0.8), Inches(4.3), Inches(11.8), Inches(2.6), '', 16, DARK_TEXT)
tf = box.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]
p.text = '对当代的启示'; p.font.size = Pt(24); p.font.color.rgb = NAVY; p.font.bold = True; p.font.name = 'Microsoft YaHei'

add_p(tf, '', 8, DARK_TEXT)
add_p(tf, '• 一个失败的民族在战后认真思过，幡然变计，是对殉国者最大的尊崇、最好的纪念', 18, DARK_TEXT)
add_p(tf, '• 失败并不可怕，关键在于能否从失败中汲取教训', 18, DARK_TEXT)
add_p(tf, '• 开放可能带来短期阵痛，但长期来看是摆脱循环的新途径', 18, DARK_TEXT)
add_p(tf, '• 历史学家应当具备远距离的思辨力', 18, DARK_TEXT)
add_p(tf, '• 日本的明治维新证明：失败的民族仍有机会再度辉煌', 18, ACCENT_RED, True)

# ====================================================================
# SLIDE 18: 结语
# ====================================================================
s = prs.slides.add_slide(prs.slide_layouts[6])
bg_navy(s)

# 装饰边框
rect(s, Inches(0.4), Inches(0.4), Inches(12.5), Inches(6.7), DARK_NAVY)
rect(s, Inches(0.5), Inches(0.5), Inches(12.3), Inches(6.5), DEEP_NAVY)
rect(s, Inches(1.2), Inches(1.2), Inches(0.08), Inches(5.1), GOLD)
rect(s, Inches(1.35), Inches(1.4), Inches(0.03), Inches(4.7), BRIGHT_GOLD)

# 标题
tb(s, Inches(1.8), Inches(1.8), Inches(9), Inches(1.0), '天朝的崩溃', 52, CREAM, True)
tb(s, Inches(1.8), Inches(2.9), Inches(9), Inches(0.7), '鸦片战争再研究', 34, GOLD, False)
rect(s, Inches(1.8), Inches(3.7), Inches(6), Inches(0.03), GOLD)
tb(s, Inches(1.8), Inches(4.0), Inches(9), Inches(0.5), '茅海建  著', 22, WARM_CREAM)
tb(s, Inches(1.8), Inches(4.7), Inches(9), Inches(0.5), '生活·读书·新知三联书店', 14, GRAY)

# 右下角装饰
diamond(s, Inches(10.5), Inches(5.5), Inches(1.2), Inches(1.2), RGBColor(0x25, 0x38, 0x58))
diamond(s, Inches(10.7), Inches(5.7), Inches(0.8), Inches(0.8), RGBColor(0x2F, 0x45, 0x68))

# 底部引文
card(s, Inches(1.8), Inches(5.3), Inches(8), Inches(1.2), RGBColor(0x15, 0x22, 0x3D))
box = tb(s, Inches(2.0), Inches(5.4), Inches(7.6), Inches(1.0), '', 14, GOLD)
tf = box.text_frame; tf.word_wrap = True; tf.paragraphs[0].alignment = PP_ALIGN.CENTER
p = tf.paragraphs[0]
p.text = '"中国的现代化一日未完成，鸦片战争的意义就一分不会减。"'
p.font.size = Pt(16); p.font.color.rgb = GOLD; p.font.italic = True; p.font.name = 'Microsoft YaHei'; p.alignment = PP_ALIGN.CENTER

# ===== 保存 =====
output = '/root/.openclaw/JDI-agent/天朝的崩溃_精美版.pptx'
prs.save(output)
print(f'Done! {len(prs.slides)} slides saved to {output}')
